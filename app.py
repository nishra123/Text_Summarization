import os
import streamlit as st
from transformers import pipeline
from docx import Document
import PyPDF2
from pptx import Presentation
import pandas as pd

def generate_summary(file_type, uploaded_file, chunk_size=1000):
    if file_type == 'email':
        # Load the JSON file into a pandas DataFrame
        uploaded_content = uploaded_file.getvalue().decode('utf-8')
        df = pd.read_json(uploaded_content)

        # Function to generate summary for a specific email
        def generate_email_summary(email_index):
            # Ensure that the email index is within the range of emails
            if email_index < 0 or email_index >= len(df):
                st.error("Invalid email index.")
                return

            # Instantiate a text summarization pipeline
            summarizer = pipeline('summarization')

            # Get the content of the email at the specified index
            email_content = df['body'].iloc[email_index]

            # Generate a summary using the text summarization pipeline
            summary = summarizer(email_content, min_length=len(email_content), do_sample=False)

            # Print the original email content and the generated summary
            st.subheader("Original Email Content:")
            st.write(email_content)
            st.subheader("Generated Summary:")
            st.write(summary[0]['summary_text'])

        return generate_email_summary

    elif file_type == 'docx':
        # Function to generate summary for a specific docx file
        def generate_docx_summary():
            # Load the docx file
            doc = Document(uploaded_file)

            # Instantiate a text summarization pipeline
            summarizer = pipeline('summarization')

            # Initialize an empty string to store the document text
            doc_text = ""

            # Iterate through paragraphs in the document
            for para in doc.paragraphs:
                doc_text += para.text + "\n"

            # Split the document text into smaller chunks
            chunks = [doc_text[i:i+chunk_size] for i in range(0, len(doc_text), chunk_size)]

            # Generate summaries for each chunk
            for i, chunk in enumerate(chunks):
                # Generate a summary using the text summarization pipeline
                summary = summarizer(chunk, min_length=50, max_length=150, do_sample=False)

                # Print the generated summary for the chunk
                st.write(summary[0]['summary_text'])

        return generate_docx_summary

    elif file_type == 'pdf':
        # Function to generate summary for a specific PDF file
        def generate_pdf_summary():
            # Open the PDF file
            pdf_reader = PyPDF2.PdfReader(uploaded_file)

            # Initialize an empty string to store the PDF text
            pdf_text = ""

            # Extract text from each page of the PDF file
            for page_num in range(len(pdf_reader.pages)):
                pdf_text += pdf_reader.pages[page_num].extract_text()

            # Split the PDF text into smaller chunks
            chunks = [pdf_text[i:i + chunk_size] for i in range(0, len(pdf_text), chunk_size)]

            # Instantiate a text summarization pipeline
            summarizer = pipeline('summarization')

            # Generate summaries for each chunk
            for i, chunk in enumerate(chunks):
                # Generate a summary using the text summarization pipeline
                summary = summarizer(chunk, min_length=50, max_length=150, do_sample=False)

                # Print the generated summary for the chunk
                st.write(summary[0]['summary_text'])

        return generate_pdf_summary

    elif file_type == 'pptx':
        # Function to generate summary for a specific PowerPoint (PPT) file
        def generate_ppt_summary():
            # Load the PowerPoint presentation
            ppt = Presentation(uploaded_file)

            # Instantiate a text summarization pipeline
            summarizer = pipeline('summarization')

            # Initialize an empty string to store the PowerPoint text
            ppt_text = ""

            # Extract text from each slide of the PowerPoint presentation
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        ppt_text += shape.text + "\n"

            # Split the PowerPoint text into smaller chunks
            chunks = [ppt_text[i:i + chunk_size] for i in range(0, len(ppt_text), chunk_size)]

            # Generate summaries for each chunk
            for i, chunk in enumerate(chunks):
                # Generate a summary using the text summarization pipeline
                summary = summarizer(chunk, min_length=50, max_length=150, do_sample=False)

                # Print the generated summary for the chunk
                st.write(summary[0]['summary_text'])

        return generate_ppt_summary

    elif file_type in ['xls', 'xlsx']:
        # Function to generate summary for a specific Excel file
        def generate_excel_summary():
            # Load the Excel file into a pandas DataFrame
            df = pd.read_excel(uploaded_file)

            # Concatenate all the cells in the Excel sheet
            excel_text = df.stack().str.strip().str.cat(sep='\n')

            # Split the Excel text into smaller chunks
            chunks = [excel_text[i:i + chunk_size] for i in range(0, len(excel_text), chunk_size)]

            # Instantiate a text summarization pipeline
            summarizer = pipeline('summarization')

            # Generate summaries for each chunk
            for i, chunk in enumerate(chunks):
                # Generate a summary using the text summarization pipeline
                summary = summarizer(chunk, min_length=50, max_length=150, do_sample=False)

                # Print the generated summary for the chunk
                st.write(summary[0]['summary_text'])

        return generate_excel_summary

    elif file_type == 'txt':
        # Function to generate summary for a specific text file
        def generate_txt_summary():
            # Read the text file
            text_content = uploaded_file.getvalue().decode('utf-8')

            # Split the text content into smaller chunks
            chunks = [text_content[i:i + chunk_size] for i in range(0, len(text_content), chunk_size)]

            # Instantiate a text summarization pipeline
            summarizer = pipeline('summarization')

            # Generate summaries for each chunk
            for i, chunk in enumerate(chunks):
                # Generate a summary using the text summarization pipeline
                summary = summarizer(chunk, min_length=50, max_length=150, do_sample=False)

                # Print the generated summary for the chunk
                st.write(summary[0]['summary_text'])

        return generate_txt_summary

    else:
        st.error("Unsupported file type.")

def main():
    st.title("File Summary Generator")

    # File upload section
    st.sidebar.header("Upload File")
    uploaded_file = st.sidebar.file_uploader("Choose a file", type=["json", "pdf", "docx", "xls", "xlsx", "pptx", "txt"])

    if uploaded_file is not None:
        file_type = uploaded_file.name.split('.')[-1]

        # Process the uploaded file based on its type
        if file_type == 'json':
            st.sidebar.write("JSON file detected.")
            st.sidebar.info("Email file detected. Please select an email index to generate the summary.")

            # Load the JSON file into a pandas DataFrame
            df = pd.read_json(uploaded_file.getvalue().decode('utf-8'))

            # Select email index
            email_index = st.sidebar.number_input("Email Index", value=0, min_value=0, max_value=len(df)-1)

            if st.sidebar.button("Generate Summary"):
                generate_email_summary = generate_summary('email', uploaded_file)
                generate_email_summary(email_index)

        elif file_type == 'pdf':
            st.sidebar.write("PDF file detected.")
            st.sidebar.info("PDF file detected. Click the button below to generate the summary.")

            if st.sidebar.button("Generate Summary"):
                generate_pdf_summary = generate_summary('pdf', uploaded_file)
                generate_pdf_summary()

        elif file_type == 'docx':
            st.sidebar.write("DOCX file detected.")
            st.sidebar.info("DOCX file detected. Click the button below to generate the summary.")

            if st.sidebar.button("Generate Summary"):
                generate_docx_summary = generate_summary('docx', uploaded_file)
                generate_docx_summary()

        elif file_type == 'pptx':
            st.sidebar.write("PPTX file detected.")
            st.sidebar.info("PowerPoint (PPTX) file detected. Click the button below to generate the summary.")

            if st.sidebar.button("Generate Summary"):
                generate_ppt_summary = generate_summary('pptx', uploaded_file)
                generate_ppt_summary()

        elif file_type in ['xls', 'xlsx']:
            st.sidebar.write("Excel file detected.")
            st.sidebar.info("Excel file detected. Click the button below to generate the summary.")

            if st.sidebar.button("Generate Summary"):
                generate_excel_summary = generate_summary('excel', uploaded_file)
                generate_excel_summary()

        elif file_type == 'txt':
            st.sidebar.write("Text file detected.")
            st.sidebar.info("Text file detected. Click the button below to generate the summary.")

            if st.sidebar.button("Generate Summary"):
                generate_txt_summary = generate_summary('txt', uploaded_file)
                generate_txt_summary()

        else:
            st.sidebar.error("Unsupported file format. Please upload a valid JSON, PDF, DOCX, Excel, PowerPoint, or Text file.")

if __name__ == "__main__":
    main()
